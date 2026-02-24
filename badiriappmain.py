import streamlit as st
import pandas as pd
from datetime import datetime
import os
import io
import requests
import base64
import json
import sqlite3

# Try to load the PowerPoint library securely
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# Try to load Plotly for the Calendar Timeline
try:
    import plotly.express as px
    HAS_PLOTLY = True
except ImportError:
    HAS_PLOTLY = False

# --- 1. APP CONFIGURATION ---
st.set_page_config(page_title="Marumo Technologies - Badiri App", layout="wide")

DB_NAME = "badiri_backend.db"
os.makedirs("attachments", exist_ok=True) # Ensure attachment folder exists

# --- 2. POWERPOINT GENERATOR ---
def create_ppt(df, sub_df):
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Badiri App Status Report"
    title_slide.placeholders[1].text = f"Marumo Technologies\nGenerated on {datetime.now().strftime('%Y-%m-%d')}"
    
    metrics_slide = prs.slides.add_slide(prs.slide_layouts[1])
    metrics_slide.shapes.title.text = "Executive Summary"
    tf = metrics_slide.shapes.placeholders[1].text_frame
    total = len(df)
    completed = len(df[df["Status"] == "Completed"])
    pending = len(df[df["Status"] == "Pending"])
    tf.text = f"Total Main Tasks: {total}"
    p1 = tf.add_paragraph()
    p1.text = f"✅ Completed Tasks: {completed}"
    p2 = tf.add_paragraph()
    p2.text = f"⏳ Pending Tasks: {pending}"
    
    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream

# --- 3. DATABASE ENGINE ---
def init_db_migration():
    conn = sqlite3.connect(DB_NAME)
    
    csv_mapping = {
        "tasks": "badiri_db.csv",
        "subtasks": "badiri_subtasks.csv",
        "users": "badiri_users.csv",
        "chat": "badiri_chat.csv",
        "mail": "badiri_mail.csv"
    }
    
    for table_name, csv_file in csv_mapping.items():
        if os.path.exists(csv_file):
            try:
                check = pd.read_sql(f"SELECT count(*) FROM {table_name}", conn)
                if check.iloc[0,0] > 0: continue 
            except: pass 
            
            df = pd.read_csv(csv_file)
            df.to_sql(table_name, conn, if_exists="replace", index=False)
            os.rename(csv_file, f"{csv_file}.backup")
    conn.close()

init_db_migration()

def load_data(table_name, default_columns):
    conn = sqlite3.connect(DB_NAME)
    try:
        df = pd.read_sql(f"SELECT * FROM {table_name}", conn)
    except:
        df = pd.DataFrame(columns=default_columns)
        df.to_sql(table_name, conn, if_exists="replace", index=False)
    conn.close()
    
    for col in default_columns:
        if col not in df.columns:
            if col == "Status": df[col] = "Active" if table_name == "users" else "Pending"
            elif col == "Role": df[col] = "Standard"
            elif col == "Password": df[col] = "1234" 
            elif col == "Read": df[col] = "No" 
            else: df[col] = ""
    return df

def save_data(df, table_name):
    conn = sqlite3.connect(DB_NAME)
    df.to_sql(table_name, conn, if_exists="replace", index=False)
    conn.close()

def show_inline_msg(location):
    if "inline_msg" in st.session_state and st.session_state.inline_msg.get("loc") == location:
        st.success(st.session_state.inline_msg["msg"])
        st.session_state.inline_msg = {} 

if "task_db" not in st.session_state: st.session_state.task_db = load_data("tasks", ["Project", "Task Name", "Assignee", "Status", "Date Added", "Due Date", "Comments", "Attachments"])
if "subtask_db" not in st.session_state: st.session_state.subtask_db = load_data("subtasks", ["Project", "Parent Task", "Subtask Name", "Assignee", "Status", "Date Added", "Due Date", "Comments", "Attachments"])
if "user_db" not in st.session_state: st.session_state.user_db = load_data("users", ["Full Name", "Email", "Phone Number", "Status", "Role", "Password"])
if "chat_db" not in st.session_state: st.session_state.chat_db = load_data("chat", ["Timestamp", "User", "Message"])
if "mail_db" not in st.session_state: st.session_state.mail_db = load_data("mail", ["Timestamp", "From", "To", "Subject", "Message", "Read"])
if "ai_suggestions" not in st.session_state: st.session_state.ai_suggestions = []
if "chat_ai_suggestions" not in st.session_state: st.session_state.chat_ai_suggestions = [] 
if "plan_ai_suggestions" not in st.session_state: st.session_state.plan_ai_suggestions = [] 
if "inline_msg" not in st.session_state: st.session_state.inline_msg = {}

if "logged_in" not in st.session_state:
    st.session_state.logged_in = False
    st.session_state.current_user = ""
    st.session_state.user_role = "Standard"
    st.session_state.is_admin = False

active_users = st.session_state.user_db[st.session_state.user_db["Status"] == "Active"] if not st.session_state.user_db.empty else pd.DataFrame()
user_list = active_users["Full Name"].tolist() if not active_users.empty else ["Unassigned"]

# --- 4. MAIN APP ROUTING ---
if not st.session_state.logged_in:
    st.title("🔒 Login to Badiri App")
    st.markdown("Welcome to the Marumo Technologies workspace.")
    with st.form("login_form"):
        email_input = st.text_input("Email Address")
        pass_input = st.text_input("Password", type="password")
        if st.form_submit_button("Login"):
            if email_input.strip().lower() == "admin" and pass_input == "Admin123":
                st.session_state.logged_in = True
                st.session_state.current_user = "Master Admin"
                st.session_state.user_role = "Admin"
                st.session_state.is_admin = True
                st.session_state.inline_msg = {"loc": "top", "msg": "✅ Logged in successfully as Master Admin!"}
                st.rerun()
            else:
                safe_db = st.session_state.user_db.copy()
                safe_db["Email"] = safe_db["Email"].astype(str).str.strip().str.lower()
                safe_db["Password"] = safe_db["Password"].astype(str).str.strip()
                user_match = safe_db[(safe_db["Email"] == email_input.strip().lower()) & (safe_db["Password"] == pass_input.strip()) & (safe_db["Status"] == "Active")]
                if not user_match.empty:
                    idx = user_match.index[0]
                    st.session_state.logged_in = True
                    st.session_state.current_user = st.session_state.user_db.at[idx, "Full Name"]
                    st.session_state.user_role = st.session_state.user_db.at[idx, "Role"]
                    st.session_state.is_admin = (st.session_state.user_role == "Admin")
                    st.session_state.inline_msg = {"loc": "top", "msg": f"✅ Welcome back, {st.session_state.current_user}!"}
                    st.rerun()
                else:
                    st.error("❌ Invalid Credentials")

else:
    with st.sidebar:
        st.header("Badiri App")
        st.caption(f"User: {st.session_state.current_user}")
        
        unread_count = len(st.session_state.mail_db[(st.session_state.mail_db["To"] == st.session_state.current_user) & (st.session_state.mail_db["Read"] == "No")])
        if unread_count > 0:
            st.error(f"📬 {unread_count} Unread Mail(s)")
            
        if st.button("🚪 Logout"):
            st.session_state.logged_in = False
            st.rerun()
        st.divider()
        if st.session_state.is_admin:
            st.subheader("👤 Register User")
            show_inline_msg("sidebar_admin") 
            with st.form("user_form", clear_on_submit=True):
                u_n = st.text_input("Name")
                u_e = st.text_input("Email")
                u_r = st.selectbox("Role", ["Standard", "Admin", "Viewer Only"])
                u_p = st.text_input("Password", type="password")
                if st.form_submit_button("Create User"):
                    new_u = pd.DataFrame([{"Full Name": u_n, "Email": u_e, "Phone Number": "", "Status": "Active", "Role": u_r, "Password": u_p}])
                    st.session_state.user_db = pd.concat([st.session_state.user_db, new_u], ignore_index=True)
                    save_data(st.session_state.user_db, "users")
                    st.session_state.inline_msg = {"loc": "sidebar_admin", "msg": f"✅ New user '{u_n}' created!"}
                    st.rerun()

    st.title("🛠️ Project Management Dashboard")
    show_inline_msg("top") 
    
    # --- STICKY NAVIGATION MENU ---
    nav_options = ["🏠 My Desk"]
    if st.session_state.user_role != "Viewer Only": nav_options.append("📁 Project Workspace")
    nav_options.append("📅 Project Calendar")
    nav_options.append("📊 Reports")
    nav_options.append("💬 Team Communications")
    if st.session_state.user_role != "Viewer Only": nav_options.append("🧠 AI Project Manager")
    if st.session_state.is_admin: nav_options.append("🛡️ Admin Console")
    
    active_tab = st.radio("Main Menu", nav_options, horizontal=True, label_visibility="collapsed", key="main_nav")
    st.divider()

    df = st.session_state.task_db
    sub_df_all = st.session_state.subtask_db

    # ==========================================
    # --- TAB 1: MY DESK ---
    # ==========================================
    if active_tab == "🏠 My Desk":
        st.subheader(f"👋 Welcome, {st.session_state.current_user}!")
        st.write("") 
        
        my_main = df[(df["Assignee"] == st.session_state.current_user) & (df["Status"] != "Completed")]
        my_sub = sub_df_all[(sub_df_all["Assignee"] == st.session_state.current_user) & (sub_df_all["Status"] != "Completed")]
        
        inbox_tasks = []
        active_tasks = []
        
        for real_idx, row in my_main.iterrows():
            is_unacknowledged = (row['Status'] == "Pending" and st.session_state.current_user not in str(row['Comments']))
            t_data = {"Type": "Main", "Idx": real_idx, "Project": row["Project"], "Name": row["Task Name"], "Status": row["Status"], "Due": row["Due Date"], "Comments": str(row["Comments"]), "Attachments": str(row.get("Attachments", ""))}
            if is_unacknowledged: inbox_tasks.append(t_data)
            else: active_tasks.append(t_data)
            
        for real_idx, row in my_sub.iterrows():
            is_unacknowledged = (row['Status'] == "Pending" and st.session_state.current_user not in str(row['Comments']))
            t_data = {"Type": "Sub", "Idx": real_idx, "Project": row["Project"], "Name": row["Subtask Name"], "Status": row["Status"], "Due": row["Due Date"], "Comments": str(row["Comments"]), "Attachments": str(row.get("Attachments", ""))}
            if is_unacknowledged: inbox_tasks.append(t_data)
            else: active_tasks.append(t_data)

        # --- INBOX SECTION ---
        st.markdown("### ⚡ Inbox: Action Required")
        show_inline_msg("desk_inbox") 
        st.caption("These are new assignments. Open them to Accept the work or Revert to someone else.")
        
        if len(inbox_tasks) == 0:
            st.info("✅ Inbox Zero! You have no new tasks waiting.")
        else:
            for t in inbox_tasks:
                with st.expander(f"🔴 NEW: {t['Project']} - {t['Name']} (Due: {t['Due']})"):
                    st.write(f"**Current Notes:** {t['Comments'] if pd.notna(t['Comments']) and t['Comments'].strip() and t['Comments'] != 'nan' else 'No notes provided.'}")
                    with st.form(f"inbox_form_{t['Type']}_{t['Idx']}"):
                        action = st.radio("Action:", ["✅ Accept Task (Move to Workspace)", "↩️ Revert Task (Reassign)"], horizontal=True)
                        c1, c2 = st.columns(2)
                        revert_user = c1.selectbox("If reverting, send to:", user_list)
                        notes = c2.text_input("Add a comment / reason:")
                        
                        if st.form_submit_button("Confirm Action"):
                            timestamp = datetime.now().strftime('%Y-%m-%d %H:%M')
                            base_cmt = t['Comments'] if pd.notna(t['Comments']) and t['Comments'].strip() != "nan" else ""
                            
                            if "Accept" in action:
                                note_text = notes.strip() if notes.strip() else "Task formally accepted."
                                new_cmt = base_cmt + f"\n[{timestamp}] {st.session_state.current_user} ACCEPTED: {note_text}"
                                if t['Type'] == "Main":
                                    st.session_state.task_db.at[t['Idx'], "Status"] = "In Progress"
                                    st.session_state.task_db.at[t['Idx'], "Comments"] = new_cmt
                                    save_data(st.session_state.task_db, "tasks")
                                else:
                                    st.session_state.subtask_db.at[t['Idx'], "Status"] = "In Progress"
                                    st.session_state.subtask_db.at[t['Idx'], "Comments"] = new_cmt
                                    save_data(st.session_state.subtask_db, "subtasks")
                                
                                st.session_state.inline_msg = {"loc": "desk_inbox", "msg": f"✅ Task '{t['Name']}' Accepted and moved to your active workspace!"}
                                st.rerun()
                            else:
                                note_text = notes.strip() if notes.strip() else "Task reverted."
                                new_cmt = base_cmt + f"\n[{timestamp}] {st.session_state.current_user} REVERTED to {revert_user}: {note_text}"
                                if t['Type'] == "Main":
                                    st.session_state.task_db.at[t['Idx'], "Assignee"] = revert_user
                                    st.session_state.task_db.at[t['Idx'], "Comments"] = new_cmt
                                    save_data(st.session_state.task_db, "tasks")
                                else:
                                    st.session_state.subtask_db.at[t['Idx'], "Assignee"] = revert_user
                                    st.session_state.subtask_db.at[t['Idx'], "Comments"] = new_cmt
                                    save_data(st.session_state.subtask_db, "subtasks")
                                    
                                st.session_state.inline_msg = {"loc": "desk_inbox", "msg": f"✅ Task Reverted and reassigned to {revert_user}!"}
                                st.rerun()

        st.divider()
        
        # --- ACTIVE PENDING TASKS ---
        st.markdown("### 📌 My Pending Tasks")
        show_inline_msg("desk_active") 
        st.caption("Click on any task below to log your progress, add receipts/files, or complete it.")
        
        if len(active_tasks) == 0:
            st.info("You don't have any active tasks currently in progress.")
        else:
            for t in active_tasks:
                icon = "⏳" if t['Status'] == "Pending" else "🚀"
                with st.expander(f"{icon} [{t['Type']}] {t['Project']} - {t['Name']} ({t['Status']})"):
                    st.write(f"**Due Date:** {t['Due']}")
                    st.write(f"**Current Notes:**\n{t['Comments'] if pd.notna(t['Comments']) and t['Comments'].strip() and t['Comments'] != 'nan' else 'No notes provided.'}")
                    
                    if pd.notna(t['Attachments']) and t['Attachments'] != "" and t['Attachments'] != "nan":
                        st.markdown("**📎 Task Attachments:**")
                        att_files = t['Attachments'].split('|')
                        for file_path in att_files:
                            if os.path.exists(file_path):
                                with open(file_path, "rb") as f:
                                    st.download_button(label=f"⬇️ Download {os.path.basename(file_path).split('_', 1)[-1]}", data=f, file_name=os.path.basename(file_path).split('_', 1)[-1], key=f"dl_{t['Type']}_{t['Idx']}_{file_path}")
                    
                    st.write("")
                    with st.form(f"update_active_{t['Type']}_{t['Idx']}"):
                        new_status = st.selectbox("Update Status", ["Pending", "In Progress", "Completed"], index=["Pending", "In Progress", "Completed"].index(t['Status']))
                        added_comment = st.text_area("Add a progress update / final notes:")
                        
                        uploaded_file = st.file_uploader("Upload Document / Receipt (Optional)")
                        
                        if st.form_submit_button("💾 Save Progress"):
                            final_comments = t['Comments'] if pd.notna(t['Comments']) and t['Comments'].strip() != "nan" else ""
                            if added_comment.strip():
                                timestamp = datetime.now().strftime('%Y-%m-%d %H:%M')
                                final_comments = final_comments.strip() + f"\n[{timestamp}] {st.session_state.current_user}: {added_comment.strip()}"
                                
                            final_atts = t['Attachments'] if pd.notna(t['Attachments']) and t['Attachments'] != "nan" else ""
                            if uploaded_file is not None:
                                safe_filename = f"{datetime.now().strftime('%Y%m%d%H%M%S')}_{uploaded_file.name.replace('|', '')}"
                                file_path = os.path.join("attachments", safe_filename)
                                with open(file_path, "wb") as f:
                                    f.write(uploaded_file.getbuffer())
                                final_atts = file_path if not final_atts else final_atts + "|" + file_path
                                
                            if t['Type'] == "Main":
                                st.session_state.task_db.at[t['Idx'], "Status"] = new_status
                                st.session_state.task_db.at[t['Idx'], "Comments"] = final_comments
                                st.session_state.task_db.at[t['Idx'], "Attachments"] = final_atts
                                save_data(st.session_state.task_db, "tasks")
                            else:
                                st.session_state.subtask_db.at[t['Idx'], "Status"] = new_status
                                st.session_state.subtask_db.at[t['Idx'], "Comments"] = final_comments
                                st.session_state.subtask_db.at[t['Idx'], "Attachments"] = final_atts
                                save_data(st.session_state.subtask_db, "subtasks")
                                
                            st.session_state.inline_msg = {"loc": "desk_active", "msg": f"✅ Progress saved for '{t['Name']}'! Status: {new_status}"}
                            st.rerun()

                    if t['Type'] == "Main":
                        st.markdown("---")
                        st.markdown("##### ➕ Create Subtask")
                        with st.form(f"quick_add_sub_{t['Idx']}"):
                            s_name = st.text_input("Subtask Name")
                            c1, c2 = st.columns(2)
                            default_user_idx = user_list.index(st.session_state.current_user) if st.session_state.current_user in user_list else 0
                            s_assignee = c1.selectbox("Assign To", user_list, index=default_user_idx)
                            s_due = c2.date_input("Due Date")
                            
                            if st.form_submit_button("Create Subtask"):
                                if s_name:
                                    new_sub = pd.DataFrame([{
                                        "Project": t['Project'], 
                                        "Parent Task": t['Name'], 
                                        "Subtask Name": s_name, 
                                        "Assignee": s_assignee, 
                                        "Status": "Pending", 
                                        "Date Added": datetime.now().strftime("%Y-%m-%d"), 
                                        "Due Date": str(s_due), 
                                        "Comments": "",
                                        "Attachments": ""
                                    }])
                                    st.session_state.subtask_db = pd.concat([st.session_state.subtask_db, new_sub], ignore_index=True)
                                    save_data(st.session_state.subtask_db, "subtasks")
                                    st.session_state.inline_msg = {"loc": "desk_active", "msg": f"✅ Subtask '{s_name}' created under '{t['Name']}'!"}
                                    st.rerun()
                                else:
                                    st.error("Please provide a subtask name.")

    # ==========================================
    # --- TAB 2: WORKSPACE ---
    # ==========================================
    elif active_tab == "📁 Project Workspace":
        st.subheader("📁 Project Workspace")
        
        existing_projects = df["Project"].unique().tolist() if not df.empty else []
        c1, c2 = st.columns([1, 2])
        project_selection = c1.selectbox("Select Workspace", ["-- Choose a Project --", "✨ Create New Project"] + existing_projects, key="ws_proj_sel")
        active_project = c2.text_input("Enter New Project Name", placeholder="e.g. Leririma Games 2026") if project_selection == "✨ Create New Project" else (project_selection if project_selection != "-- Choose a Project --" else None)

        if active_project:
            st.divider()
            st.markdown(f"### 📂 Project: {active_project}")
            
            proj_df = df[df["Project"] == active_project]
            proj_sub_df = sub_df_all[sub_df_all["Project"] == active_project]
            
            m1, m2, m3 = st.columns(3)
            tot_tasks = len(proj_df)
            comp_tasks = len(proj_df[proj_df["Status"] == "Completed"])
            pct = (comp_tasks / tot_tasks) if tot_tasks > 0 else 0.0
            
            m1.metric("Total Main Tasks", tot_tasks)
            m2.metric("Subtasks Attached", len(proj_sub_df))
            m3.metric("Overall Completion", f"{int(pct*100)}%")
            st.progress(pct)
            st.write("")
            
            # STICKY TABS for Workspace Sub-Menu
            pw_tab = st.radio("Workspace Operations", ["🗂️ Project Board", "➕ Add New Task", "⚙️ Edit Tasks & Subtasks"], horizontal=True, label_visibility="collapsed", key="pw_nav")
            st.write("")
            
            if pw_tab == "🗂️ Project Board":
                if proj_df.empty:
                    st.info("No tasks in this project yet. Go to 'Add New Task' to get started.")
                else:
                    for real_idx, m_row in proj_df.iterrows():
                        with st.container(border=True):
                            icon = "✅" if m_row['Status'] == "Completed" else "🔹"
                            st.markdown(f"#### {icon} {m_row['Task Name']}")
                            st.caption(f"**Assignee:** {m_row['Assignee']} | **Status:** {m_row['Status']} | **Due:** {m_row['Due Date']}")
                            
                            if pd.notna(m_row['Comments']) and str(m_row['Comments']).strip() != "nan" and str(m_row['Comments']).strip() != "":
                                st.write(f"**Notes:** {m_row['Comments']}")
                                
                            m_subs = proj_sub_df[proj_sub_df["Parent Task"] == m_row["Task Name"]]
                            if not m_subs.empty:
                                st.markdown("**Subtasks:**")
                                st.dataframe(m_subs[["Subtask Name", "Assignee", "Status", "Due Date"]], hide_index=True, use_container_width=True)
            
            elif pw_tab == "➕ Add New Task":
                st.markdown("#### Create a New Main Task")
                show_inline_msg("ws_add_main") 
                with st.form("workspace_add_task_form", clear_on_submit=True):
                    t_name = st.text_input("Task Name")
                    t_assignee = st.selectbox("Assign To", user_list)
                    t_status = st.selectbox("Status", ["Pending", "In Progress", "Completed"])
                    t_due = st.date_input("Due Date")
                    t_comments = st.text_area("Comments")
                    if st.form_submit_button("Add Task") and t_name:
                        new_task = pd.DataFrame([{"Project": active_project, "Task Name": t_name, "Assignee": t_assignee, "Status": t_status, "Date Added": datetime.now().strftime("%Y-%m-%d"), "Due Date": str(t_due), "Comments": t_comments, "Attachments": ""}])
                        st.session_state.task_db = pd.concat([st.session_state.task_db, new_task], ignore_index=True)
                        save_data(st.session_state.task_db, "tasks")
                        st.session_state.inline_msg = {"loc": "ws_add_main", "msg": f"✅ New task '{t_name}' added to {active_project}!"}
                        st.rerun()

            elif pw_tab == "⚙️ Edit Tasks & Subtasks":
                update_col1, update_col2 = st.columns(2)
                
                with update_col1:
                    st.markdown("**✏️ Edit Main Task**")
                    show_inline_msg("ws_upd_main")
                    if not proj_df.empty:
                        task_dict = {idx: row["Task Name"] for idx, row in proj_df.iterrows()}
                        selected_idx = st.selectbox("Select Task to Edit", options=list(task_dict.keys()), format_func=lambda x: task_dict[x])
                        if selected_idx is not None:
                            curr_assig = df.at[selected_idx, "Assignee"]
                            with st.form("workspace_update_form"):
                                new_assignee = st.selectbox("Reassign To", user_list, index=user_list.index(curr_assig) if curr_assig in user_list else 0)
                                new_status = st.selectbox("Status", ["Pending", "In Progress", "Completed"], index=["Pending", "In Progress", "Completed"].index(df.at[selected_idx, "Status"]))
                                new_comments = st.text_area("Comments", value=str(df.at[selected_idx, "Comments"]))
                                if st.form_submit_button("Save Updates"):
                                    if new_assignee != curr_assig: 
                                        new_comments += f"\n[Forwarded to {new_assignee}]"
                                    st.session_state.task_db.at[selected_idx, "Assignee"] = new_assignee
                                    st.session_state.task_db.at[selected_idx, "Status"] = new_status
                                    st.session_state.task_db.at[selected_idx, "Comments"] = new_comments
                                    save_data(st.session_state.task_db, "tasks")
                                    st.session_state.inline_msg = {"loc": "ws_upd_main", "msg": "✅ Task successfully updated!"}
                                    st.rerun()
                                    
                with update_col2:
                    st.markdown("**⚙️ Manage Subtasks**")
                    show_inline_msg("ws_sub_mng")
                    if not proj_df.empty:
                        parent_task = st.selectbox("Select Parent Task:", ["-- Select --"] + proj_df["Task Name"].tolist())
                        if parent_task != "-- Select --":
                            with st.expander("➕ Add Subtask", expanded=False):
                                with st.form("add_sub_form", clear_on_submit=True):
                                    s_name = st.text_input("Subtask Name")
                                    s_assignee = st.selectbox("Assign To", user_list)
                                    s_due = st.date_input("Due Date")
                                    if st.form_submit_button("Create Subtask") and s_name:
                                        new_sub = pd.DataFrame([{"Project": active_project, "Parent Task": parent_task, "Subtask Name": s_name, "Assignee": s_assignee, "Status": "Pending", "Date Added": datetime.now().strftime("%Y-%m-%d"), "Due Date": str(s_due), "Comments": "", "Attachments": ""}])
                                        st.session_state.subtask_db = pd.concat([st.session_state.subtask_db, new_sub], ignore_index=True)
                                        save_data(st.session_state.subtask_db, "subtasks")
                                        st.session_state.inline_msg = {"loc": "ws_sub_mng", "msg": f"✅ New subtask '{s_name}' added!"}
                                        st.rerun()
                                        
                            active_subtasks = sub_df_all[(sub_df_all["Project"] == active_project) & (sub_df_all["Parent Task"] == parent_task)]
                            if not active_subtasks.empty:
                                with st.expander("✏️ Edit Subtask", expanded=False):
                                    sub_dict = {idx: row["Subtask Name"] for idx, row in active_subtasks.iterrows()}
                                    sub_idx = st.selectbox("Select Subtask", options=list(sub_dict.keys()), format_func=lambda x: sub_dict[x])
                                    if sub_idx is not None:
                                        with st.form("update_sub_form"):
                                            s_curr_status = sub_df_all.at[sub_idx, "Status"]
                                            new_s_status = st.selectbox("Status", ["Pending", "In Progress", "Completed"], index=["Pending", "In Progress", "Completed"].index(s_curr_status))
                                            if st.form_submit_button("Save Subtask Updates"):
                                                st.session_state.subtask_db.at[sub_idx, "Status"] = new_s_status
                                                save_data(st.session_state.subtask_db, "subtasks")
                                                st.session_state.inline_msg = {"loc": "ws_sub_mng", "msg": "✅ Subtask successfully updated!"}
                                                st.rerun()

    # ==========================================
    # --- TAB 3: PROJECT CALENDAR ---
    # ==========================================
    elif active_tab == "📅 Project Calendar":
        st.subheader("📅 Project Calendar & Visual Timeline")
        st.markdown("Track exactly when tasks begin and when they are due.")
        
        m_cal = df.copy()
        s_cal = sub_df_all.copy()
        
        if m_cal.empty and s_cal.empty:
            st.info("No tasks to display on the calendar.")
        else:
            if not m_cal.empty: m_cal["Task Display"] = "[Main] " + m_cal["Task Name"]
            if not s_cal.empty: s_cal["Task Display"] = "[Sub] " + s_cal["Subtask Name"]
            
            cal_df = pd.concat([
                m_cal[["Project", "Task Display", "Assignee", "Status", "Date Added", "Due Date"]] if not m_cal.empty else pd.DataFrame(columns=["Project", "Task Display", "Assignee", "Status", "Date Added", "Due Date"]),
                s_cal[["Project", "Task Display", "Assignee", "Status", "Date Added", "Due Date"]] if not s_cal.empty else pd.DataFrame(columns=["Project", "Task Display", "Assignee", "Status", "Date Added", "Due Date"])
            ], ignore_index=True)
            
            cal_df["Start"] = pd.to_datetime(cal_df["Date Added"], errors='coerce')
            cal_df["End"] = pd.to_datetime(cal_df["Due Date"], errors='coerce')
            cal_df["End"] = cal_df.apply(lambda x: x["Start"] + pd.Timedelta(days=1) if pd.isna(x["End"]) or x["Start"] == x["End"] else x["End"], axis=1)
            cal_df = cal_df.dropna(subset=["Start", "End"])
            
            if not cal_df.empty:
                cal_df = cal_df.sort_values("End")
                if HAS_PLOTLY:
                    fig = px.timeline(cal_df, x_start="Start", x_end="End", y="Task Display", color="Project", hover_name="Assignee", hover_data=["Status", "Due Date"], height=500)
                    fig.update_yaxes(autorange="reversed") 
                    st.plotly_chart(fig, use_container_width=True)
                else:
                    st.warning("⚠️ **Plotly Required:** To see the visual timeline, please add `plotly` to your requirements.txt file and reboot.")

                st.divider()
                st.markdown("#### 🚨 Upcoming Deadlines (Next 7 Days)")
                today = pd.Timestamp.now().normalize()
                next_week = today + pd.Timedelta(days=7)
                
                upcoming = cal_df[(cal_df["End"] >= today) & (cal_df["End"] <= next_week) & (cal_df["Status"] != "Completed")]
                if upcoming.empty:
                    st.success("You are all clear! No pending deadlines in the next 7 days. 🎉")
                else:
                    st.dataframe(upcoming[["Project", "Task Display", "Assignee", "Status", "Due Date"]].rename(columns={"Task Display": "Task"}), hide_index=True, use_container_width=True)

    # ==========================================
    # --- TAB 4: REPORTS ---
    # ==========================================
    elif active_tab == "📊 Reports":
        if df.empty:
            st.info("No tasks to report on.")
        else:
            c1, c2, c3 = st.columns(3)
            c1.metric("Total Main Tasks", len(df))
            c2.metric("✅ Tasks Completed", len(df[df["Status"] == "Completed"]))
            c3.metric("Total Subtasks", len(sub_df_all))
            
            st.divider()
            st.subheader("📊 Project Health Dashboard")
            for proj in df["Project"].unique():
                p_df = df[df["Project"] == proj]
                p_tot = len(p_df)
                p_comp = len(p_df[p_df["Status"] == "Completed"])
                p_pct = (p_comp / p_tot) if p_tot > 0 else 0.0
                st.write(f"**{proj}** ({p_comp}/{p_tot} Tasks Completed)")
                st.progress(p_pct)
                st.write("")
            
            st.divider()
            st.subheader("📈 Team Performance & Capacity Matrix")
            all_assignments = pd.concat([df[["Assignee", "Status"]], sub_df_all[["Assignee", "Status"]]], ignore_index=True)
            
            matrix_data = []
            for user in all_assignments["Assignee"].dropna().unique():
                u_tasks = all_assignments[all_assignments["Assignee"] == user]
                u_tot = len(u_tasks)
                u_comp = len(u_tasks[u_tasks["Status"] == "Completed"])
                u_pct = int((u_comp / u_tot) * 100) if u_tot > 0 else 0
                matrix_data.append({"Team Member": user, "Total Load": u_tot, "Completed": u_comp, "Efficiency %": u_pct})
            
            if matrix_data:
                matrix_df = pd.DataFrame(matrix_data)
                st.dataframe(
                    matrix_df,
                    column_config={
                        "Efficiency %": st.column_config.ProgressColumn("Efficiency Rate", format="%d%%", min_value=0, max_value=100)
                    },
                    hide_index=True, use_container_width=True
                )
            
            st.divider()
            st.subheader("📥 Export Center")
            ex1, ex2 = st.columns(2)
            with ex1:
                if HAS_PPTX:
                    st.download_button("📊 Download PowerPoint", data=create_ppt(df, sub_df_all), file_name=f"Report_{datetime.now().strftime('%Y%m%d')}.pptx")
            with ex2:
                st.download_button("📈 Download CSV Export", data=df.to_csv(index=False).encode('utf-8'), file_name=f"Data_{datetime.now().strftime('%Y%m%d')}.csv")

    # ==========================================
    # --- TAB 5: TEAM COMMUNICATIONS ---
    # ==========================================
    elif active_tab == "💬 Team Communications":
        st.subheader("💬 Team Communications")
        st.markdown("Chat with your team in real-time or send formal, secure internal mail.")
        
        # Sticky navigation for Comm Tabs
        comm_tab = st.radio("Communication Actions", ["💬 Global Team Chat", "📥 Mail Inbox", "📤 Compose Mail"], horizontal=True, label_visibility="collapsed", key="comm_nav")
        st.write("")
        
        if comm_tab == "💬 Global Team Chat":
            st.session_state.chat_db = load_data("chat", ["Timestamp", "User", "Message"])
            chat_container = st.container(height=400)
            with chat_container:
                if st.session_state.chat_db.empty:
                    st.caption("No messages yet. Say hello!")
                else:
                    for _, msg in st.session_state.chat_db.tail(20).iterrows():
                        is_me = (msg["User"] == st.session_state.current_user)
                        with st.chat_message("user" if is_me else "assistant"):
                            st.markdown(f"**{msg['User']}** <span style='font-size:0.8em; color:gray;'>({msg['Timestamp']})</span>", unsafe_allow_html=True)
                            st.write(msg["Message"])
            
            with st.form("chat_form", clear_on_submit=True):
                m = st.text_input("Type your message to the team...")
                c1, c2 = st.columns(2)
                if c1.form_submit_button("📨 Send Message") and m:
                    new_c = pd.DataFrame([{"Timestamp": datetime.now().strftime("%H:%M"), "User": st.session_state.current_user, "Message": m}])
                    st.session_state.chat_db = pd.concat([st.session_state.chat_db, new_c], ignore_index=True)
                    save_data(st.session_state.chat_db, "chat")
                    st.rerun()
                if c2.form_submit_button("🔄 Refresh Chat"): st.rerun()

        elif comm_tab == "📥 Mail Inbox":
            show_inline_msg("mail_inbox")
            my_mail = st.session_state.mail_db[st.session_state.mail_db["To"] == st.session_state.current_user]
            if my_mail.empty:
                st.info("Your inbox is empty.")
            else:
                for idx, row in my_mail.sort_index(ascending=False).iterrows():
                    unread_tag = "🔴 [NEW]" if row["Read"] == "No" else "⚪"
                    with st.expander(f"{unread_tag} {row['Subject']} - From: {row['From']} ({row['Timestamp']})"):
                        st.write(row["Message"])
                        if row["Read"] == "No":
                            if st.button("Mark as Read", key=f"read_mail_{idx}"):
                                st.session_state.mail_db.at[idx, "Read"] = "Yes"
                                save_data(st.session_state.mail_db, "mail")
                                st.session_state.inline_msg = {"loc": "mail_inbox", "msg": "✅ Mail marked as read."}
                                st.rerun()

        elif comm_tab == "📤 Compose Mail":
            show_inline_msg("mail_compose") 
            with st.form("compose_mail_form", clear_on_submit=True):
                to_user = st.selectbox("Send To:", user_list)
                subject = st.text_input("Subject")
                msg = st.text_area("Your Message")
                if st.form_submit_button("Send Secure Mail"):
                    if subject and msg:
                        new_mail = pd.DataFrame([{
                            "Timestamp": datetime.now().strftime("%Y-%m-%d %H:%M"),
                            "From": st.session_state.current_user,
                            "To": to_user,
                            "Subject": subject,
                            "Message": msg,
                            "Read": "No"
                        }])
                        st.session_state.mail_db = pd.concat([st.session_state.mail_db, new_mail], ignore_index=True)
                        save_data(st.session_state.mail_db, "mail")
                        st.session_state.inline_msg = {"loc": "mail_compose", "msg": f"✅ Secure mail successfully sent to {to_user}!"}
                        st.rerun()
                    else:
                        st.error("Please fill in subject and message.")

    # ==========================================
    # --- TAB 6: AI PROJECT MANAGER ---
    # ==========================================
    elif active_tab == "🧠 AI Project Manager":
        st.subheader("🧠 Your AI Project Manager")
        st.markdown("👋 **Hello! I am your Badiri AI Assistant.**\n\nI can help you automate your workspace. Generate a whole project from a single sentence, scan meeting minutes, or mine the Team Chat!")
        
        gemini_key = st.text_input("🔑 Gemini API Key", type="password")
        st.divider()
        
        # --- AUTO-PLANNER ---
        st.markdown("#### 🪄 Auto-Generate Project Plan")
        st.caption("Tell me what you are trying to do, and I will build a full step-by-step project plan.")
        
        with st.form("ai_planner_form"):
            plan_prompt = st.text_input("Describe the project:", placeholder="e.g., Organize a constituency cleanup campaign for 500 people...")
            if st.form_submit_button("🏗️ Build Project Plan"):
                if gemini_key and plan_prompt:
                    with st.spinner("Drafting tasks..."):
                        full_prompt = f"You are an expert Project Manager. I need to plan this: '{plan_prompt}'. Create a comprehensive project plan. Return strictly a JSON list of objects with these keys: 'Project' (a short unifying project name), 'Task Name', 'Assignee' (always output 'Unassigned'). Do not include markdown or explanations."
                        res = requests.post(f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent?key={gemini_key}", json={"contents": [{"parts": [{"text": full_prompt}]}]}).json()
                        if 'candidates' in res: 
                            st.session_state.plan_ai_suggestions = json.loads(res['candidates'][0]['content']['parts'][0]['text'].replace("```json","").replace("```","").strip())
        
        if st.session_state.plan_ai_suggestions and st.session_state.is_admin:
            show_inline_msg("ai_plan") 
            with st.form("plan_approval"):
                st.write("**Select tasks to import into Workspace:**")
                plan_sels = [st.checkbox(f"{it['Task Name']}", value=True, key=f"plan_c_{i}") for i, it in enumerate(st.session_state.plan_ai_suggestions)]
                if st.form_submit_button("✅ Approve Selected Plan"):
                    added = 0
                    for i, sel in enumerate(plan_sels):
                        if sel: 
                            st.session_state.task_db = pd.concat([st.session_state.task_db, pd.DataFrame([{"Project": st.session_state.plan_ai_suggestions[i]['Project'], "Task Name": st.session_state.plan_ai_suggestions[i]['Task Name'], "Assignee": st.session_state.plan_ai_suggestions[i]['Assignee'], "Status": "Pending", "Date Added": datetime.now().strftime("%Y-%m-%d"), "Due Date": datetime.now().strftime("%Y-%m-%d"), "Comments": "AI Auto-Generated Plan", "Attachments": ""}])], ignore_index=True)
                            added += 1
                    save_data(st.session_state.task_db, "tasks")
                    st.session_state.plan_ai_suggestions = []
                    st.session_state.inline_msg = {"loc": "ai_plan", "msg": f"✅ {added} tasks imported from the AI Planner!"}
                    st.rerun()

        st.divider()

        # --- EXTRACT FROM IMAGE ---
        st.markdown("#### 📷 Extract Tasks from Minutes (Image)")
        img_file = st.file_uploader("Upload Minutes", type=["jpg", "png"])
        if st.button("🔍 Analyze Minutes"):
            if gemini_key and img_file:
                with st.spinner("Analyzing document..."):
                    b64 = base64.b64encode(img_file.read()).decode('utf-8')
                    prompt = f"Extract tasks as JSON list with keys: Project, Task Name, Assignee. Use only these names: {user_list}"
                    res = requests.post(f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent?key={gemini_key}", json={"contents": [{"parts": [{"text": prompt}, {"inline_data": {"mime_type": "image/jpeg", "data": b64}}]}]}).json()
                    if 'candidates' in res: st.session_state.ai_suggestions = json.loads(res['candidates'][0]['content']['parts'][0]['text'].replace("```json","").replace("```","").strip())
        
        if st.session_state.ai_suggestions and st.session_state.is_admin:
            show_inline_msg("ai_img") 
            with st.form("img_approval"):
                st.write("**Select items to import into Workspace:**")
                img_sels = [st.checkbox(f"{it['Project']} | {it['Task Name']} ({it['Assignee']})", value=True, key=f"img_c_{i}") for i, it in enumerate(st.session_state.ai_suggestions)]
                if st.form_submit_button("✅ Approve Selected"):
                    added = 0
                    for i, sel in enumerate(img_sels):
                        if sel: 
                            st.session_state.task_db = pd.concat([st.session_state.task_db, pd.DataFrame([{"Project": st.session_state.ai_suggestions[i]['Project'], "Task Name": st.session_state.ai_suggestions[i]['Task Name'], "Assignee": st.session_state.ai_suggestions[i]['Assignee'], "Status": "Pending", "Date Added": datetime.now().strftime("%Y-%m-%d"), "Due Date": datetime.now().strftime("%Y-%m-%d"), "Comments": "AI extracted", "Attachments": ""}])], ignore_index=True)
                            added += 1
                    save_data(st.session_state.task_db, "tasks")
                    st.session_state.ai_suggestions = []
                    st.session_state.inline_msg = {"loc": "ai_img", "msg": f"✅ {added} task(s) imported from the document!"}
                    st.rerun()

        st.divider()

        # --- EXTRACT FROM CHAT ---
        st.markdown("#### 💬 Extract Tasks from Chat logs")
        if st.button("🧠 Analyze Chat Logs"):
            if gemini_key and not st.session_state.chat_db.empty:
                with st.spinner("Mining chat..."):
                    transcript = "\n".join([f"{r['User']}: {r['Message']}" for _, r in st.session_state.chat_db.tail(30).iterrows()])
                    prompt = f"Extract tasks from chat as JSON list with keys: Project, Task Name, Assignee. Names: {user_list}\n\nCHAT:\n{transcript}"
                    res = requests.post(f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent?key={gemini_key}", json={"contents": [{"parts": [{"text": prompt}]}]}).json()
                    if 'candidates' in res: st.session_state.chat_ai_suggestions = json.loads(res['candidates'][0]['content']['parts'][0]['text'].replace("```json","").replace("```","").strip())

        if st.session_state.chat_ai_suggestions and st.session_state.is_admin:
            show_inline_msg("ai_chat") 
            with st.form("chat_approval"):
                st.write("**Select chat promises to import:**")
                chat_sels = [st.checkbox(f"{it['Project']} | {it['Task Name']} ({it['Assignee']})", value=True, key=f"chat_c_{i}") for i, it in enumerate(st.session_state.chat_ai_suggestions)]
                if st.form_submit_button("✅ Approve Selected"):
                    added = 0
                    for i, sel in enumerate(chat_sels):
                        if sel: 
                            st.session_state.task_db = pd.concat([st.session_state.task_db, pd.DataFrame([{"Project": st.session_state.chat_ai_suggestions[i]['Project'], "Task Name": st.session_state.chat_ai_suggestions[i]['Task Name'], "Assignee": st.session_state.chat_ai_suggestions[i]['Assignee'], "Status": "Pending", "Date Added": datetime.now().strftime("%Y-%m-%d"), "Due Date": datetime.now().strftime("%Y-%m-%d"), "Comments": "Chat AI extracted", "Attachments": ""}])], ignore_index=True)
                            added += 1
                    save_data(st.session_state.task_db, "tasks")
                    st.session_state.chat_ai_suggestions = []
                    st.session_state.inline_msg = {"loc": "ai_chat", "msg": f"✅ {added} task(s) automatically extracted from chat!"}
                    st.rerun()

    # ==========================================
    # --- TAB 7: ADMIN ---
    # ==========================================
    elif active_tab == "🛡️ Admin Console" and st.session_state.is_admin:
        st.subheader("🛡️ Admin Console")
        
        st.markdown("#### 👥 User Management")
        if not st.session_state.user_db.empty: st.dataframe(st.session_state.user_db, hide_index=True, use_container_width=True)
        
        show_inline_msg("admin_edit") 
        user_to_update = st.selectbox("Select User to Edit", ["-- Select User --"] + st.session_state.user_db["Full Name"].tolist())
        if user_to_update != "-- Select User --":
            curr_user = st.session_state.user_db[st.session_state.user_db["Full Name"] == user_to_update].iloc[0]
            idx = st.session_state.user_db.index[st.session_state.user_db["Full Name"] == user_to_update].tolist()[0]
            with st.form("update_user_details"):
                c1, c2 = st.columns(2)
                n_n = c1.text_input("Name", value=curr_user["Full Name"])
                n_e = c2.text_input("Email", value=curr_user["Email"])
                n_p = c1.text_input("Phone", value=str(curr_user["Phone Number"]).replace('nan',''))
                n_s = c2.selectbox("Status", ["Active", "Suspended", "Blocked"], index=["Active", "Suspended", "Blocked"].index(curr_user["Status"]))
                n_r = c1.selectbox("Role", ["Standard", "Admin", "Viewer Only"], index=["Standard", "Admin", "Viewer Only"].index(curr_user["Role"]))
                n_pw = c2.text_input("Password", value=curr_user["Password"], type="password")
                if st.form_submit_button("Save Changes"):
                    st.session_state.user_db.loc[idx] = [n_n, n_e, n_p, n_s, n_r, n_pw]
                    save_data(st.session_state.user_db, "users")
                    st.session_state.inline_msg = {"loc": "admin_edit", "msg": f"✅ Profile for {n_n} updated successfully."}
                    st.rerun()

# --- END OF FILE ---
